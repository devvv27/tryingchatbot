from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt


DEPLOYMENT_LINK = "https://tryingchatbot-fqskya4z4ncu5xdzhr973r.streamlit.app"
GITHUB_LINK = "https://github.com/devvv27/tryingchatbot"
OUTPUT_FILE = "AI_Booking_Assistant_Presentation_v2.pptx"

PRIMARY = RGBColor(15, 50, 78)
SECONDARY = RGBColor(41, 128, 185)
ACCENT = RGBColor(19, 177, 154)
SURFACE = RGBColor(237, 245, 250)
TEXT_DARK = RGBColor(30, 35, 40)


def style_title(shape):
    p = shape.text_frame.paragraphs[0]
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = PRIMARY


def add_footer(slide, text="AI Booking Assistant | NeoStats AI Engineer Use Case"):
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.0),
        Inches(7.1),
        Inches(13.33),
        Inches(0.4),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = SURFACE
    line.line.fill.background()
    tf = line.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = PRIMARY


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    style_title(slide.shapes.title)
    sub = slide.placeholders[1].text_frame.paragraphs[0]
    sub.font.size = Pt(20)
    sub.font.color.rgb = TEXT_DARK

    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.6),
        Inches(5.5),
        Inches(12.1),
        Inches(1.2),
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = SURFACE
    banner.line.color.rgb = SECONDARY
    t = banner.text_frame
    t.clear()
    p = t.paragraphs[0]
    p.text = "Chat + RAG + Booking Workflow + Admin Dashboard + Email Confirmation"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    add_footer(slide)


def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    style_title(slide.shapes.title)
    text_frame = slide.shapes.placeholders[1].text_frame
    text_frame.clear()

    for i, item in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0

        p.text = text
        p.level = level
        p.font.size = Pt(22 if level == 0 else 18)
        p.font.color.rgb = TEXT_DARK

    add_footer(slide)


def add_agenda_slide(prs):
    add_bullet_slide(
        prs,
        "Agenda",
        [
            "Problem statement and goals",
            "System design and architecture",
            "RAG pipeline and booking flow",
            "Admin dashboard and validations",
            "Deployment, links, and future enhancements",
        ],
    )


def add_kpi_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Implementation Coverage"
    style_title(slide.shapes.title)

    cards = [
        ("6", "Core modules implemented"),
        ("4", "Tools: RAG, DB, Email, Web Search"),
        ("25", "Short-term memory window"),
        ("100%", "Mandatory requirements covered"),
    ]
    left = 0.8
    for number, label in cards:
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(2.0),
            Inches(2.95),
            Inches(2.8),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = SURFACE
        card.line.color.rgb = SECONDARY

        tf = card.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = number
        p1.font.size = Pt(34)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY

        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(17)
        p2.font.color.rgb = TEXT_DARK
        left += 3.15

    add_footer(slide)


def add_two_column_slide(prs, title, left_title, left_points, right_title, right_points):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    style_title(slide.shapes.title)

    left_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.3))
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.3))

    left_tf = left_box.text_frame
    left_tf.word_wrap = True
    p = left_tf.paragraphs[0]
    p.text = left_title
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = PRIMARY
    for point in left_points:
        q = left_tf.add_paragraph()
        q.text = point
        q.level = 1
        q.font.size = Pt(18)
        q.font.color.rgb = TEXT_DARK

    right_tf = right_box.text_frame
    right_tf.word_wrap = True
    p2 = right_tf.paragraphs[0]
    p2.text = right_title
    p2.font.bold = True
    p2.font.size = Pt(24)
    p2.font.color.rgb = PRIMARY
    for point in right_points:
        q2 = right_tf.add_paragraph()
        q2.text = point
        q2.level = 1
        q2.font.size = Pt(18)
        q2.font.color.rgb = TEXT_DARK

    add_footer(slide)


def add_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "System Architecture"
    style_title(slide.shapes.title)

    boxes = [
        ("User (Streamlit UI)", 0.7),
        ("Chat Logic\nIntent + Memory + Routing", 2.0),
        ("RAG Pipeline\nPDF -> Chunk -> Embed -> Retrieve", 3.3),
        ("Booking Tools\nDB Save + Email", 4.6),
        ("SQLite + Admin Dashboard", 5.9),
    ]

    for i, (text, top) in enumerate(boxes):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(1.0),
            Inches(top),
            Inches(11.2),
            Inches(0.9),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = SURFACE if i % 2 == 0 else RGBColor(227, 238, 247)
        shape.line.color.rgb = SECONDARY
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = PRIMARY

        if i < len(boxes) - 1:
            connector = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.DOWN_ARROW,
                Inches(6.35),
                Inches(top + 0.82),
                Inches(0.5),
                Inches(0.45),
            )
            connector.fill.solid()
            connector.fill.fore_color.rgb = ACCENT
            connector.line.color.rgb = ACCENT

    add_footer(slide)


def add_booking_flow_visual(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Booking Flow (Visual)"
    style_title(slide.shapes.title)

    steps = [
        "Intent Detection",
        "Extract Known Fields",
        "Ask Missing Fields",
        "Show Summary",
        "Get Confirmation",
        "Save + Email + Booking ID",
    ]

    x = 0.6
    y = 2.2
    for idx, step in enumerate(steps, start=1):
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(2.0),
            Inches(1.4),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = SURFACE
        box.line.color.rgb = SECONDARY
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = f"{idx}. {step}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = PRIMARY

        if idx < len(steps):
            arrow = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
                Inches(x + 1.95),
                Inches(y + 0.45),
                Inches(0.45),
                Inches(0.45),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT
            arrow.line.color.rgb = ACCENT

        x += 2.1

    add_footer(slide)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "AI Booking Assistant",
        "Use Case Submission Presentation\nCandidate: Dewang Choudhary",
    )

    add_agenda_slide(prs)

    add_kpi_slide(prs)

    add_bullet_slide(
        prs,
        "Use Case",
        [
            "Build a chat-based AI assistant for bookings.",
            "Support PDF-based RAG for information answering.",
            "Detect booking intent and collect details in multi-turn chat.",
            "Confirm details before storage and send confirmation email.",
            "Provide an Admin Dashboard to view and filter bookings.",
        ],
    )

    add_bullet_slide(
        prs,
        "Solution Overview and Approach",
        [
            "Frontend: Streamlit with chat interface, PDF upload, and admin page.",
            "LLM: Gemini (primary) with provider-configured setup.",
            "RAG: PDF text extraction, chunking, TF-IDF embeddings, top-k retrieval.",
            "Booking flow: intent detection, slot filling, validation, confirmation step.",
            "Persistence: SQLite schema for customers and bookings.",
            "Notifications: SMTP email tool with graceful failure handling.",
        ],
    )

    add_architecture_slide(prs)

    add_bullet_slide(
        prs,
        "RAG Design",
        [
            "User uploads one or more PDFs in Streamlit.",
            "System extracts text from PDF pages.",
            "Text is split into overlapping chunks.",
            "Chunks are embedded into a lightweight vector representation.",
            "For each query, top relevant chunks are retrieved and passed to LLM.",
            "Fallback retrieval handles broad summary-style questions.",
        ],
    )

    add_bullet_slide(
        prs,
        "Booking Flow",
        [
            "1. Detect booking intent from user input.",
            "2. Extract known fields: name, email, phone, type, date, time.",
            "3. Ask only missing fields with validation.",
            "4. Summarize collected details.",
            "5. Ask explicit confirmation from user.",
            "6. On confirm: save booking in DB and trigger email.",
            "7. Return booking ID and maintain recent chat memory.",
        ],
    )

    add_booking_flow_visual(prs)

    add_two_column_slide(
        prs,
        "Admin Dashboard and UX",
        "Admin Dashboard",
        [
            "View all stored bookings",
            "Filter by name, email, date",
            "See booking status and timestamps",
        ],
        "User Experience",
        [
            "Clear user vs assistant chat messages",
            "PDF upload and processing status",
            "Friendly runtime and validation messages",
        ],
    )

    add_bullet_slide(
        prs,
        "Validation and Error Handling",
        [
            "Handles invalid or non-extractable PDFs.",
            "Validates email, date (YYYY-MM-DD), and time (HH:MM).",
            "Catches DB insert/connection failures.",
            "Handles email delivery failures gracefully.",
            "Shows user-friendly messages and status logs.",
        ],
    )

    add_bullet_slide(
        prs,
        "Challenges and Future Improvements",
        [
            "Challenge: scanned PDFs may have low extractable text.",
            "Future: add OCR fallback for image-only PDFs.",
            "Future: upgrade to semantic vector DB for larger corpora.",
            "Future: add booking edit/cancel/export actions in admin panel.",
            "Future: add STT/TTS and richer UI animations.",
        ],
    )

    add_bullet_slide(
        prs,
        "Submission Links",
        [
            f"Deployment Link: {DEPLOYMENT_LINK}",
            f"GitHub Repository: {GITHUB_LINK}",
            "App includes chat, RAG, booking flow, DB storage, email tool, and admin dashboard.",
        ],
    )

    prs.save(OUTPUT_FILE)
    print(f"Created: {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
